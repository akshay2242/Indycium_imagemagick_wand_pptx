from wand.image import Image
from pptx import Presentation
from pptx.util import Inches,Cm,Pt
from pptx.text.text import Font
import os

pr1 = Presentation()

def create_ppt(logo_image_url,list_of_images):
    # Inserting watermark on images
    with Image(filename=logo_image_url) as logo:
        logo.resize(1600,600)                     
        for i in range(len(list_of_images)):
                with Image(filename=list_of_images[i]) as img:
                    img.composite(image=logo, left=0, top=0)
                    img.save(filename=f'output_images_with_logo/new_image{i+1}_with_logo.jpg') # Saving output images
                                
    for i in range(1,6):
        # Creating slide in ppt
        slide_register = pr1.slide_layouts[1]
        slide = pr1.slides.add_slide(slide_register)
        
        # Creating title of slide
        title1 = slide.shapes.title
        title1.text = f"Sample Title {i}"
        title1.text_frame.paragraphs[0].font.name = 'Arial'
        title1.text_frame.paragraphs[0].font.size = Pt(40)
        title1.width = Inches(5)
        title1.height = Inches(2)
        title1.left = Cm(3)

        # Creating subtitle of slide
        subtitle = slide.placeholders[1]
        subtitle.text = f"Sample Subtitle {i}"
        subtitle.width = Inches(5)
        subtitle.height = Inches(0.5)
        subtitle.left = Cm(1)
        subtitle.top = Cm(3.5)

        # Inserting picture in slide
        picture = f"output_images_with_logo/new_image{i}_with_logo.jpg"
        from_left = Cm(2)
        from_top = Cm(5)
        add_picture = slide.shapes.add_picture(picture, from_left, from_top, width=Inches(4), height=Inches(5))

# To check directory is present or not 
list_dirs = os.listdir(os.getcwd())
if 'output_images_with_logo' in list_dirs:
    pass
else:
    os.mkdir('output_images_with_logo')

logo_image_url = 'input_images/nike_black.png' # logo image
list_of_images = ['input_images/image1.jpg', 'input_images/image2.jpg', 'input_images/image3.jpg', 'input_images/image4.jpg', 'input_images/image5.jpg',] # List of input images

create_ppt(logo_image_url,list_of_images) 

pr1.save('myppt.pptx')